import streamlit as st
from pptx import Presentation
from pptx.util import Inches
import io

# 初始化 Streamlit 應用程式
st.title("旅遊行程編排工具")

# 側欄選項
st.sidebar.title("新增旅遊日期頁面")
if st.sidebar.button("新增日期頁面"):
    st.session_state.dates = st.session_state.get("dates", []) + [{"date": "", "activities": []}]

# 顯示日期頁面
if "dates" not in st.session_state:
    st.session_state.dates = []

for i, date_page in enumerate(st.session_state.dates):
    st.subheader(f"日期頁面 {i+1}")
    date_page["date"] = st.text_input(f"日期 {i+1}", date_page.get("date", ""))
    st.write("活動：")
    for j, activity in enumerate(date_page.get("activities", [])):
        activity["name"] = st.text_input(f"活動 {j+1} 名稱", activity.get("name", ""))
        activity["description"] = st.text_area(f"活動 {j+1} 描述", activity.get("description", ""))
    if st.button(f"新增活動 {i+1}"):
        date_page.setdefault("activities", []).append({"name": "", "description": ""})

# 上傳 PPT 模板文件
uploaded_file = st.file_uploader("選擇 PPT 模板文件（可選）", type=["pptx"])

# 輸出行程為 PPT
if st.button("生成 PPT"):
    if uploaded_file is not None:
        # 載入上傳的 PPT 模板文件
        prs = Presentation(uploaded_file)
    else:
        # 生成一個新的 PPT 文件
        prs = Presentation()

    # 添加旅遊行程內容
    for date_page in st.session_state.dates:
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        title.text = f"日期: {date_page['date']}"
        txBox = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(5))
        tf = txBox.text_frame
        for activity in date_page.get("activities", []):
            p = tf.add_paragraph()
            p.text = f"活動名稱: {activity['name']}\n描述: {activity['description']}\n"

    # 將 PPT 保存到內存中
    ppt_io = io.BytesIO()
    prs.save(ppt_io)
    ppt_io.seek(0)

    # 提供下載連結
    st.download_button(
        label="下載 PPT",
        data=ppt_io,
        file_name="旅遊行程.pptx",
        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
    )
